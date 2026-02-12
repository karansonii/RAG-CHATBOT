# src/data_loader.py
import os
from pathlib import Path
from docx import Document
from pptx import Presentation
import PyPDF2
from paddleocr import PaddleOCR
import logging
from src.utils import ConfigManager
import io
import hashlib
import json
import time
import numpy as np
from PIL import Image
from tqdm import tqdm
# Setup logging
logger = logging.getLogger(__name__)
# Load configuration
rag_config = ConfigManager.get_rag_config()
CHUNK_SIZE = rag_config['chunk_size']
CHUNK_OVERLAP = rag_config['overlap']
PPT_CHUNK_SIZE = rag_config['ppt_chunk_size']
PPT_CHUNK_OVERLAP = rag_config['ppt_overlap']
# OCR configuration
ocr_lang = ConfigManager.get('OCR_LANGUAGE', 'en')
ocr_angle = ConfigManager.get('OCR_ANGLE_DETECTION', True, config_type=bool)
# Initialize OCR reader once
ocr_reader = PaddleOCR(use_angle_cls=ocr_angle, lang=ocr_lang)
logger.info(f"OCR initialized: lang={ocr_lang}, angle_detection={ocr_angle}")
# ================== CONFIG FROM OCR.PY ==================
log_level = ConfigManager.get('LOG_LEVEL', 'INFO')
logging.basicConfig(
    filename=ConfigManager.get('DATA_LOADER_LOG', 'data_loader.log'),
    level=getattr(logging, log_level),
    format='%(asctime)s [%(levelname)s] %(message)s'
)
CONF_THRESHOLD = ConfigManager.get('OCR_CONFIDENCE_THRESHOLD', 0.7, config_type=float)
CACHE_DIR = ConfigManager.get('CACHE_DIR', './data/cache')
os.makedirs(CACHE_DIR, exist_ok=True)
# ================== CACHE HELPERS ==================
def get_file_hash(file_path):
    hash_md5 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()
def get_slide_hash(slide):
    hasher = hashlib.sha256()
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            hasher.update(shape.text.strip().encode("utf-8"))
        if hasattr(shape, "image"):
            hasher.update(shape.image.blob)
    return hasher.hexdigest()
def get_cache_path(file_path):
    file_hash = get_file_hash(file_path)
    return os.path.join(CACHE_DIR, f"{file_hash}.json")
# ================== CENTRAL CACHE INDEX ==================
INDEX_FILE = os.path.join(CACHE_DIR, "cache_index.json")
def load_cache_index():
    if os.path.exists(INDEX_FILE):
        try:
            with open(INDEX_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logging.error(f"Error loading cache index: {e}")
    return {}
def save_cache_index(index_data):
    try:
        with open(INDEX_FILE, "w", encoding="utf-8") as f:
            json.dump(index_data, f, ensure_ascii=False, indent=4)
        logging.info(f"💾 Updated cache index at {INDEX_FILE}")
    except Exception as e:
        logging.error(f"Failed to save cache index: {e}")
def load_cached_result(file_path):
    file_name = os.path.basename(file_path)
    index = load_cache_index()
    cached_entry = index.get(file_name)
   
    if cached_entry:
        cache_file = cached_entry.get("cache_path")
        file_hash = cached_entry.get("hash")
        current_hash = get_file_hash(file_path)
        if current_hash == file_hash and os.path.exists(cache_file):
            try:
                with open(cache_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    logging.info(f"✅ Using indexed cache for {file_path}")
                    print(f"⚡ Cache found for {os.path.basename(file_path)}")
                    return data
            except Exception as e:
                logging.error(f"Error reading indexed cache: {e}")
   
    return {}
def save_cache(file_path, data):
    try:
        file_hash = get_file_hash(file_path)
        cache_file = get_cache_path(file_path) # This uses file_hash
        file_name = os.path.basename(file_path)
       
        with open(cache_file, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
        logging.info(f"💾 Cached data at {cache_file}")
        print(f"💾 Cached data at {cache_file}")
       
        index = load_cache_index()
        index[file_name] = {
            "hash": file_hash,
            "cache_path": cache_file,
            "last_updated": time.time()
        }
        save_cache_index(index)
    except Exception as e:
        logging.error(f"Failed to save cache: {e}")
        print(f"❌ Failed to save cache: {e}")
# ================== OCR CORE FROM OCR.PY ==================
def extract_images_from_slide(slide):
    images = []
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            image_stream = io.BytesIO(shape.image.blob)
            img = Image.open(image_stream)
            images.append(img)
    return images
def read_text_from_image(image):
    img_array = np.array(image)
    result = ocr_reader.ocr(img_array, cls=True)
    extracted = [
        line[1][0]
        for page in result if page is not None
        for line in page
        if line[1][1] >= CONF_THRESHOLD
    ]
    image.close()
    return "\n".join(extracted)
def process_slide(slide, slide_num):
    start_time = time.time()
    slide_text = []
    for shape in slide.shapes:
        # Extract text from shape
        if hasattr(shape, "text") and shape.text.strip():
            slide_text.append(shape.text.strip())
        # Extract tables as markdown
        if shape.shape_type == 19: # TABLE
            try:
                md_table = table_to_markdown(shape.table)
                if md_table.strip():
                    slide_text.append(md_table)
            except Exception as e:
                logging.warning(f"Failed to extract table from slide {slide_num}: {e}")
        # Extract text from image using OCR
        if hasattr(shape, "image"):
            try:
                img = Image.open(io.BytesIO(shape.image.blob))
                ocr_text = read_text_from_image(img)
                if ocr_text.strip():
                    slide_text.append(ocr_text.strip())
            except Exception as e:
                logging.error(f"OCR error on slide {slide_num}: {e}")
    elapsed = time.time() - start_time
    print(f"🕒 Slide {slide_num} processed in {elapsed:.2f}s")
    return "\n".join(slide_text)
# ================== MERGED READ_PPT_TEXT ==================
def read_ppt_text(ppt_path):
    cached_data = load_cached_result(ppt_path)
    cached_slides = cached_data.get("slides", {})
    prs = Presentation(ppt_path)
    ppt_text = {"slides": {}}
    slides_to_process = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_hash = get_slide_hash(slide)
        slide_key = f"Slide_{i}"
        cached_slide = cached_slides.get(slide_key)
        if not cached_slide or cached_slide["hash"] != slide_hash:
            slides_to_process.append((i, slide, slide_hash))
        else:
            ppt_text["slides"][slide_key] = cached_slide
    print(f"[DEBUG] Slides to OCR: {len(slides_to_process)}")
    # ✅ if cache already available
    if not slides_to_process:
        print("✅ No new or modified slides. Using cached OCR.")
        return {k: v["text"] for k, v in ppt_text["slides"].items()}
    # 🚀 Process new slides
    for slide_num, slide, slide_hash in tqdm(slides_to_process, desc="Running OCR"):
        text = process_slide(slide, slide_num)
        ppt_text["slides"][f"Slide_{slide_num}"] = {"hash": slide_hash, "text": text}
    # 💾 Save OCR cache to JSON
    save_cache(ppt_path, ppt_text)
    return {k: v["text"] for k, v in ppt_text["slides"].items()}
# ================== CACHED TEXT EXTRACTION FOR OTHER FORMATS ==================
def extract_docx_text(path):
    file_hash = get_file_hash(path)
    cached_data = load_cached_result(path)
   
    if cached_data and cached_data.get("hash") == file_hash:
        return cached_data["full_text"]
   
    try:
        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        full_text = "\n".join(paragraphs)
       
        data = {
            "hash": file_hash,
            "full_text": full_text
        }
        save_cache(path, data)
        return full_text
    except Exception as e:
        logger.error(f"Failed to extract DOCX {path}: {e}")
        return ""
def extract_pdf_text(path):
    file_hash = get_file_hash(path)
    cached_data = load_cached_result(path)
   
    if cached_data and cached_data.get("hash") == file_hash:
        return cached_data["full_text"]
   
    text = ""
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                text += page_text
       
        data = {
            "hash": file_hash,
            "full_text": text
        }
        save_cache(path, data)
        return text
    except Exception as e:
        logger.error(f"Failed to extract PDF {path}: {e}")
        return ""
# ================== CHUNKING ==================
def chunk_text(text, source_file, chunk_size=None, overlap=None, extra_meta=None):
    if chunk_size is None:
        chunk_size = CHUNK_SIZE
    if overlap is None:
        overlap = CHUNK_OVERLAP
   
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk_words = words[start:end]
        chunk_text_str = " ".join(chunk_words)
        meta = {
            "source_file": source_file,
            "start_word": start,
            "end_word": end,
            "chunk_size": len(chunk_words)
        }
        if extra_meta:
            meta.update(extra_meta)
        chunks.append({"text": chunk_text_str, "metadata": meta})
        start += chunk_size - overlap
    logger.debug(f"Created {len(chunks)} chunks from {source_file}")
    return chunks
# ================== LOAD FUNCTIONS WITH CACHING ==================
def load_docx(path):
    full_text = extract_docx_text(path)
    if not full_text:
        return []
    chunks = chunk_text(full_text, os.path.basename(path))
    logger.info(f"Loaded {len(chunks)} chunks from DOCX: {os.path.basename(path)}")
    return chunks
def load_pdf(path):
    full_text = extract_pdf_text(path)
    if not full_text:
        return []
    chunks = chunk_text(full_text, os.path.basename(path))
    logger.info(f"Loaded {len(chunks)} chunks from PDF: {os.path.basename(path)}")
    return chunks
def load_pptx(path):
    ppt_text_dict = read_ppt_text(path)
    slide_chunks = []
   
    for slide_key, text in ppt_text_dict.items():
        i = int(slide_key.split('_')[1])
        if text:
            slide_chunks.extend(
                chunk_text(
                    text,
                    os.path.basename(path),
                    chunk_size=PPT_CHUNK_SIZE,
                    overlap=PPT_CHUNK_OVERLAP,
                    extra_meta={"slide_number": i},
                )
            )
    logger.info(f"Loaded {len(slide_chunks)} chunks from PPTX: {os.path.basename(path)}")
    return slide_chunks
# ================== TABLE TO MARKDOWN ==================
def table_to_markdown(table):
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    header = " | ".join(rows[0])
    separator = " | ".join(["---"] * len(rows[0]))
    body = "\n".join([" | ".join(r) for r in rows[1:]])
    return f"{header}\n{separator}\n{body}" if body else f"{header}\n{separator}"
def load_documents_from_folder(folder_path):
    folder = Path(folder_path)
    all_chunks = []
    file_stats = {}
    index = load_cache_index()
   
    # Cleanup stale entries (files no longer in folder)
    existing_files = {f.name for f in folder.iterdir() if f.is_file()}
    stale_keys = [k for k in index if k not in existing_files]
    for k in stale_keys:
        cache_path = index[k].get("cache_path")
        if cache_path and os.path.exists(cache_path):
            os.remove(cache_path)
            logging.info(f"🗑️ Removed stale cache: {cache_path}")
        del index[k]
    if stale_keys:
        save_cache_index(index)
   
    supported_extensions = {
        '.docx': load_docx,
        '.pptx': load_pptx,
        '.pdf': load_pdf
    }
    for file in folder.iterdir():
        if not file.is_file():
            continue
        # Skip temporary or hidden files
        if file.name.startswith("~$") or file.name.startswith("."):
            logger.debug(f"Skipping temp/hidden file: {file.name}")
            continue
        ext = file.suffix.lower()
       
        if ext in supported_extensions:
            loader_func = supported_extensions[ext]
            chunks = loader_func(file)
            all_chunks.extend(chunks)
            file_stats[file.name] = len(chunks)
        else:
            logger.debug(f"Skipping unsupported file: {file.name}")
    # Log statistics
    logger.info(f"✅ Loaded {len(all_chunks)} total chunks from {len(file_stats)} files")
    for filename, count in file_stats.items():
        logger.info(f" {filename}: {count} chunks")
    return all_chunks
# Export
__all__ = [
    'load_documents_from_folder',
    'load_docx',
    'load_pdf',
    'load_pptx',
    'chunk_text',
    'read_ppt_text'
]
